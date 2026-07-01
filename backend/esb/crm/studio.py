"""GovernanceWriter + PresentationCreator — stateless content/deck generation.

Ported from coach-devon's governance/writer.py + presentation/build.py.
Both async here (unlike verifier/dossier, no ORM relationship traversal —
these have no DB dependency at all).
"""
from __future__ import annotations

import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from esb.crm import llm

# ── Governance Writer ───────────────────────────────────────────────────────

STYLE_GUIDE = """
AJ Crabill / Effective School Boards voice:
- Relentless focus on STUDENT OUTCOMES as the purpose of governance.
- Distinguish governance (the board's vision & accountability) from management
  (the superintendent's operations). Never blur them.
- Plain, confident, warm. Short sentences. No hedging, no hype.
- Respectful of board members as public servants; never condescending.
- Concrete over abstract. Name the behavior, not the buzzword.
- Avoid edu-jargon, acronyms, and corporate filler.
"""

BANNED = [
    "synergy", "leverage", "circle back", "best-in-class", "world-class",
    "cutting-edge", "paradigm", "low-hanging fruit", "move the needle",
    "at the end of the day", "going forward", "thought leader", "deep dive",
    "stakeholder buy-in", "value-add", "holistic", "robust solution",
]
_PASSIVE = re.compile(r"\b(is|are|was|were|be|been|being)\s+\w+ed\b", re.I)

WRITE_SYS = (
    "You write in the voice defined by this style guide. Produce clear, warm, "
    "student-outcomes-focused prose. Return JSON: {\"text\":\"...\"}.\n\n" + STYLE_GUIDE
)


def voice_lint(text: str) -> list[dict]:
    flags: list[dict] = []
    low = text.lower()
    for phrase in BANNED:
        if phrase in low:
            flags.append({"type": "jargon", "issue": f"banned phrase: '{phrase}'"})
    passive = _PASSIVE.findall(text)
    if len(passive) >= 3:
        flags.append({"type": "passive", "issue": f"{len(passive)} passive constructions — prefer active voice"})
    longs = [s for s in re.split(r"[.!?]", text) if len(s.split()) > 35]
    if longs:
        flags.append({"type": "length", "issue": f"{len(longs)} very long sentence(s) — tighten"})
    return flags


async def write(purpose: str, context: str = "", draft: str = "") -> dict:
    user = f"purpose: {purpose}\ncontext: {context}\n"
    if draft:
        user += f"rewrite this draft in-voice:\n{draft}"
    out = await llm.complete_json(WRITE_SYS, user)
    text = out.get("text", "")
    return {"text": text, "voice_flags": voice_lint(text)}


# ── Presentation Creator ─────────────────────────────────────────────────────

PRIMARY = RGBColor(0x12, 0x3D, 0x6B)
ACCENT = RGBColor(0xE8, 0x9B, 0x1E)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"
FOOTER = "Effective School Boards"

OUTLINE_SYS = (
    "You outline a concise, persuasive slide deck for Effective School Boards, "
    "focused on student-outcomes governance. Return JSON: "
    "{\"title\":\"...\",\"subtitle\":\"...\","
    "\"slides\":[{\"title\":\"...\",\"bullets\":[\"...\"]}]}. 5-8 slides, 3-5 bullets each."
)


async def generate_outline(topic: str) -> dict:
    return await llm.complete_json(OUTLINE_SYS, f"topic: {topic}")


def _bar(slide, prs, color, height_in):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(height_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _footer(slide, prs):
    box = slide.shapes.add_textbox(Inches(0.4), prs.slide_height - Inches(0.4),
                                    prs.slide_width - Inches(0.8), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = FOOTER
    run.font.size = Pt(9)
    run.font.color.rgb = PRIMARY
    run.font.name = FONT


def build_deck(spec: dict) -> bytes:
    """spec: {title, subtitle, slides:[{title, bullets:[...]}]} -> .pptx bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    bg.shadow.inherit = False
    _bar(s, prs, ACCENT, 0.25)
    tb = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = spec.get("title", "Untitled")
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = LIGHT
    r.font.name = FONT
    if spec.get("subtitle"):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = spec["subtitle"]
        r.font.size = Pt(20)
        r.font.color.rgb = ACCENT
        r.font.name = FONT

    for sl in spec.get("slides", []):
        s = prs.slides.add_slide(blank)
        _bar(s, prs, PRIMARY, 1.1)
        t = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12), Inches(0.8))
        r = t.text_frame.paragraphs[0].add_run()
        r.text = sl.get("title", "")
        r.font.size = Pt(26)
        r.font.bold = True
        r.font.color.rgb = LIGHT
        r.font.name = FONT

        body = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.4))
        tf = body.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(sl.get("bullets", [])):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = "•  " + str(bullet)
            r.font.size = Pt(18)
            r.font.color.rgb = DARK
            r.font.name = FONT
            p.space_after = Pt(10)
        _footer(s, prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
